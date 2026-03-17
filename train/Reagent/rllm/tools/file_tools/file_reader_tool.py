import csv
import json
import mimetypes
import os
import tempfile
import zipfile
from datetime import datetime, date
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

from rllm.tools.tool_base import Tool, ToolOutput

# Default configuration
FILE_BASE_PATH = os.getenv("FILE_BASE_PATH", "./data")
MAX_CONTENT_CHARS = int(os.getenv("FILEREADER_MAX_CHARS", "2000"))


def json_serialize_handler(obj):
    """Handle non-serializable objects for JSON encoding."""
    if isinstance(obj, (datetime, date)):
        return obj.isoformat()
    # Handle pandas Timestamp and other datetime-like objects
    if hasattr(obj, 'isoformat'):
        return obj.isoformat()
    # Handle NaN and NaT
    if str(obj) in ('nan', 'NaN', 'NaT'):
        return None
    raise TypeError(f"Object of type {type(obj)} is not JSON serializable")


# Check optional dependencies
try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except Exception:
    PANDAS_AVAILABLE = False

try:
    import xml.etree.ElementTree as ET
    XML_AVAILABLE = True
except Exception:
    XML_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PDF_AVAILABLE = True
except Exception:
    PDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except Exception:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False


class _FileReaderCore:
    """Internal utility class providing multi-format file reading support."""

    def __init__(self, max_chars: int = MAX_CONTENT_CHARS):
        self.max_chars = max_chars
        self.supported_formats = {
            ".txt": self._read_text,
            ".py": self._read_text,
            ".json": self._read_json,
            ".jsonl": self._read_jsonl,
            ".jsonld": self._read_jsonld,
            ".csv": self._read_csv,
            ".xml": self._read_xml,
            ".xlsx": self._read_excel,
            ".xls": self._read_excel,
            ".pdf": self._read_pdf,
            ".docx": self._read_docx,
            ".doc": self._read_docx,
            ".pptx": self._read_pptx,
            ".zip": self._read_zip,
            ".pdb": self._read_text,
        }

    # Public API -----------------------------------------------------------------
    def read_file(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """
        Read a file and return its content with metadata.
        
        Args:
            file_path: Path to the file. Supports ZIP internal files with format: zip_path::internal_file
            **kwargs: Additional arguments (e.g., encoding)
        
        Returns:
            Dict with 'success', 'content', 'metadata', and optionally 'error'
        """
        try:
            # Handle ZIP internal files
            if "::" in file_path:
                return self._read_zip_internal(file_path, **kwargs)

            if not os.path.exists(file_path):
                return self._error(f"File does not exist: {file_path}")

            file_ext = Path(file_path).suffix.lower()
            reader = self.supported_formats.get(file_ext)
            if not reader:
                return self._error(f"Unsupported file format: {file_ext}")

            metadata = self._get_file_info(file_path)
            content = reader(file_path, **kwargs)
            truncated = self._truncate_content(content)
            return {
                "success": True,
                "content": truncated,
                "metadata": metadata,
            }
        except Exception as exc:
            return self._error(f"Failed to read file: {str(exc)}", file_path=file_path)

    # Helpers --------------------------------------------------------------------
    def _error(self, message: str, file_path: Optional[str] = None) -> Dict[str, Any]:
        """Create an error response."""
        meta = {}
        if file_path and os.path.exists(file_path):
            meta = self._get_file_info(file_path)
        return {"success": False, "error": message, "content": None, "metadata": meta}

    def _get_file_info(self, file_path: str) -> Dict[str, Any]:
        """Extract file metadata."""
        stat = os.stat(file_path)
        return {
            "file_path": file_path,
            "file_name": os.path.basename(file_path),
            "file_size": stat.st_size,
            "file_extension": Path(file_path).suffix.lower(),
            "mime_type": mimetypes.guess_type(file_path)[0],
            "modified_time": stat.st_mtime,
        }

    def _truncate_content(self, content: Any) -> Any:
        """Truncate content if it exceeds max_chars."""
        content_str = content if isinstance(content, str) else json.dumps(
            content, ensure_ascii=False, default=json_serialize_handler
        )
        if len(content_str) <= self.max_chars:
            return content
        truncated = content_str[: self.max_chars] + f"\n... (truncated, original length: {len(content_str)} chars)"
        return truncated

    # Format readers -------------------------------------------------------------
    def _read_text(self, file_path: str, encoding: str = "utf-8", **kwargs) -> str:
        """Read plain text files with automatic encoding detection."""
        try:
            with open(file_path, "r", encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            # Try alternative encodings
            for enc in ["gbk", "gb2312", "latin-1"]:
                try:
                    with open(file_path, "r", encoding=enc) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
        raise UnicodeDecodeError("FileReader", b"", 0, 1, "Unable to decode file as text.")

    def _read_json(self, file_path: str, encoding: str = "utf-8", **kwargs) -> Union[Dict[str, Any], List[Any]]:
        """Read JSON files."""
        with open(file_path, "r", encoding=encoding) as f:
            return json.load(f)

    def _read_jsonl(self, file_path: str, encoding: str = "utf-8", **kwargs) -> List[Any]:
        """Read JSON Lines files."""
        records: List[Any] = []
        with open(file_path, "r", encoding=encoding) as f:
            for line in f:
                line = line.strip()
                if line:
                    records.append(json.loads(line))
        return records

    def _read_jsonld(self, file_path: str, encoding: str = "utf-8", **kwargs) -> Union[Dict[str, Any], List[Any]]:
        """Read JSON-LD files."""
        with open(file_path, "r", encoding=encoding) as f:
            return json.load(f)

    def _read_csv(self, file_path: str, encoding: str = "utf-8", **kwargs) -> List[Dict[str, Any]]:
        """Read CSV files, preferring pandas if available."""
        if PANDAS_AVAILABLE:
            try:
                df = pd.read_csv(file_path, encoding=encoding)
                # Convert datetime columns to string
                for col in df.columns:
                    if pd.api.types.is_datetime64_any_dtype(df[col]):
                        df[col] = df[col].astype(str)
                return df.to_dict("records")
            except Exception:
                pass
        
        # Fallback to standard csv module
        records: List[Dict[str, Any]] = []
        with open(file_path, "r", encoding=encoding, newline="") as f:
            sample = f.read(1024)
            f.seek(0)
            delimiter = csv.Sniffer().sniff(sample).delimiter
            reader = csv.DictReader(f, delimiter=delimiter)
            for row in reader:
                records.append(dict(row))
        return records

    def _read_xml(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Read XML files."""
        if not XML_AVAILABLE:
            raise RuntimeError("XML parsing is not available. Please install xml libraries.")
        
        tree = ET.parse(file_path)
        root = tree.getroot()

        def to_dict(elem):
            node: Dict[str, Any] = {}
            if elem.attrib:
                node["@attributes"] = elem.attrib
            text = (elem.text or "").strip()
            if text:
                node["#text"] = text
            for child in elem:
                child_data = to_dict(child)
                if child.tag in node:
                    if not isinstance(node[child.tag], list):
                        node[child.tag] = [node[child.tag]]
                    node[child.tag].append(child_data)
                else:
                    node[child.tag] = child_data
            return node

        return {"root_tag": root.tag, "content": to_dict(root)}

    def _read_excel(self, file_path: str, **kwargs) -> Dict[str, List[Dict[str, Any]]]:
        """Read Excel files."""
        if not PANDAS_AVAILABLE:
            raise RuntimeError("Reading Excel requires pandas + openpyxl.")
        
        excel = pd.ExcelFile(file_path)
        data: Dict[str, List[Dict[str, Any]]] = {}
        for sheet in excel.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet)
            # Convert datetime columns to string
            for col in df.columns:
                if pd.api.types.is_datetime64_any_dtype(df[col]):
                    df[col] = df[col].astype(str)
            data[sheet] = df.to_dict("records")
        return data

    def _read_pdf(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Read PDF files."""
        if not PDF_AVAILABLE:
            raise RuntimeError("Reading PDF requires PyMuPDF (fitz).")
        
        doc = fitz.open(file_path)
        pages = []
        try:
            for idx in range(doc.page_count):
                page = doc.load_page(idx)
                text = page.get_text()
                pages.append(
                    {
                        "page_number": idx + 1,
                        "text": text,
                        "text_length": len(text),
                        "width": page.rect.width,
                        "height": page.rect.height,
                    }
                )
            full_text = "\n".join(p["text"] for p in pages)
            return {
                "total_pages": len(pages),
                "pages": pages,
                "full_text": full_text,
                "metadata": doc.metadata
            }
        finally:
            doc.close()

    def _read_docx(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Read Word documents."""
        if not DOCX_AVAILABLE:
            raise RuntimeError("Reading Word documents requires python-docx.")
        
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        tables = []
        for tbl in doc.tables:
            table_rows = []
            for row in tbl.rows:
                table_rows.append([cell.text for cell in row.cells])
            tables.append(table_rows)
        return {
            "paragraphs": paragraphs,
            "tables": tables,
            "full_text": "\n".join(paragraphs)
        }

    def _read_pptx(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Read PowerPoint presentations."""
        if not PPTX_AVAILABLE:
            raise RuntimeError("Reading PowerPoint requires python-pptx.")
        
        prs = Presentation(file_path)
        slides = []
        for idx, slide in enumerate(prs.slides):
            entry = {"slide_number": idx + 1, "title": "", "content": []}
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape == getattr(slide.shapes, "title", None):
                        entry["title"] = shape.text
                    else:
                        entry["content"].append(shape.text)
            slides.append(entry)
        
        full_text = "\n".join(
            f"Slide {s['slide_number']} - {s['title']}\n" + "\n".join(s["content"])
            for s in slides
        )
        return {
            "total_slides": len(slides),
            "slides": slides,
            "full_text": full_text
        }

    def _read_zip(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Read ZIP archive and list contents with previews."""
        with zipfile.ZipFile(file_path, "r") as zf:
            file_list = zf.namelist()
            info = []
            previews = {}
            max_size = kwargs.get("max_file_size", 1024 * 1024)
            
            for name in file_list:
                zinfo = zf.getinfo(name)
                info.append(
                    {
                        "filename": name,
                        "file_size": zinfo.file_size,
                        "compressed_size": zinfo.compress_size,
                        "is_directory": name.endswith("/"),
                    }
                )
                # Preview small files
                if not name.endswith("/") and zinfo.file_size <= max_size:
                    data = zf.read(name)
                    try:
                        previews[name] = data.decode("utf-8")
                    except UnicodeDecodeError:
                        previews[name] = f"<binary file: {len(data)} bytes>"
            
            return {
                "total_files": len(file_list),
                "files_info": info,
                "files_preview": previews
            }

    def _read_zip_internal(self, path: str, **kwargs) -> Dict[str, Any]:
        """Read a file inside a ZIP archive using format: zip_path::internal_path"""
        zip_path, inner = path.split("::", 1)
        
        # 如果 zip_path 不是绝对路径，应该拼接 base_path
        if not os.path.isabs(zip_path) and hasattr(self, 'base_path') and self.base_path:
            zip_path = os.path.join(self.base_path, zip_path)
        
        if not os.path.exists(zip_path):
            return self._error(f"ZIP file does not exist: {zip_path}")
        with zipfile.ZipFile(zip_path, "r") as zf:
            if inner not in zf.namelist():
                return self._error(f"Internal file not found in ZIP: {inner}")
            
            data = zf.read(inner)
            suffix = Path(inner).suffix
            
            # Write to temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                tmp.write(data)
                tmp_path = tmp.name
            
            try:
                reader = self.supported_formats.get(suffix, self._read_text)
                content = reader(tmp_path, **kwargs)
                return {
                    "success": True,
                    "content": self._truncate_content(content),
                    "metadata": {
                        "zip_path": zip_path,
                        "internal_path": inner
                    }
                }
            finally:
                try:
                    os.unlink(tmp_path)
                except Exception:
                    pass


class FileReaderTool(Tool):
    """Tool to read local files and return text content from various file formats."""

    NAME = "file_reader"
    DESCRIPTION = "Read and extract content from various file formats including txt, json, csv, xlsx, pdf, docx, pptx, xml, zip, and more."

    def __init__(
        self,
        name: str = NAME,
        description: str = DESCRIPTION,
        max_chars: int = MAX_CONTENT_CHARS,
        base_path: str = FILE_BASE_PATH,
    ):
        """
        Initialize the FileReader tool.

        Args:
            name (str): The name of the tool.
            description (str): A description of the tool's purpose.
            max_chars (int): Maximum characters in output before truncation.
            base_path (str): Base path for relative file paths.
        """
        self.max_chars = max_chars
        self.base_path = base_path
        self.reader = _FileReaderCore(max_chars=max_chars)
        super().__init__(name=name, description=description)

    @property
    def json(self):
        """Return the tool's information in the required format."""
        return {
            "type": "function",
            "function": {
                "name": self.name,
                "description": "Read and extract content from various file formats including xml, xlsx, docs, csv, pdf, txt, json, jsonl, jsonld, pptx, py, zip, pdb.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "file_path": {
                            "type": "string",
                            "description": "The path to the local file to be read. Supports ZIP internal files with format: zip_path::internal_file_path"
                        }
                    },
                    "required": ["file_path"]
                }
            }
        }

    def forward(self, file_path: str, encoding: str = "utf-8") -> ToolOutput:
        """
        Read a file and return its content.
        
        Args:
            file_path: Path to the file. Supports ZIP internal files: zip_path::internal_file
            encoding: Text encoding (default: utf-8)
        
        Returns:
            ToolOutput containing file content or error
        """
        try:
            # Handle relative vs absolute paths
            final_path = file_path
            # 特殊处理：如果包含 "::" (ZIP 内部文件)
            if "::" in file_path:
                zip_path, internal_path = file_path.split("::", 1)
                # 如果 ZIP 路径不是绝对路径，拼接 base_path
                if self.base_path and not os.path.isabs(zip_path):
                    zip_path = os.path.join(self.base_path, zip_path)
                # 重新组合完整路径
                final_path = f"{zip_path}::{internal_path}"
            # 普通文件路径处理
            elif self.base_path and not os.path.isabs(file_path):
                final_path = os.path.join(self.base_path, file_path)

            # Read the file
            result = self.reader.read_file(final_path, encoding=encoding)
            
            if not result.get("success"):
                return ToolOutput(
                    name=self.name or "file_reader",
                    error=result.get("error", "Unknown error")
                )

            # Format the content
            content = result.get("content")
            if isinstance(content, (dict, list)):
                content_str = json.dumps(content, ensure_ascii=False, default=json_serialize_handler)
            else:
                content_str = str(content)

            return ToolOutput(
                name=self.name or "file_reader",
                output=content_str,
                metadata=result.get("metadata")
            )

        except Exception as e:
            return ToolOutput(
                name=self.name or "file_reader",
                error=f"FileReader error: {str(e)}"
            )


if __name__ == "__main__":
    # Test the tool
    tool = FileReaderTool()
    print("Testing FileReader Tool...")
    print("=" * 80)
    
    # Test with a sample file (replace with actual path)
    test_file = "./data/rl/files/rl/files/aaaxaxuuez282.zip"
    
    # if os.path.exists(test_file):
    print(f"\nReading: {test_file}")
    result = tool(file_path=test_file)
    if result.error:
        print(f"Error: {result.error}")
    else:
        print(f"Success!")
        print(f"Content preview: {result.output[:200]}...")
    # else:
    #     print(f"⚠️  Test file not found: {test_file}")
    #     print("Please set a valid file path to test")
    
    print("=" * 80)

